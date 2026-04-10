"""
File Processor — Multi-Format Ingestion Logic
=============================================

Architecture Rationale:
-----------------------
This module acts as the 'Gateway' for data entering the system. It handles the 
messy reality of file parsing so the KnowledgeBase can focus on pure text math.

Design Pattern: Strategy
The `process_single_file` function uses a file-extension matching strategy 
to route different data formats (PDF, CSV, Office Docs) to their respective
specialized libraries.

Key Techniques:
1. **Encoding Resilience**: Implements fallbacks (UTF-8 -> Latin1 -> CP1252) 
   to handle common Windows/Excel encoding issues.
2. **Page-Awareness**: Preserves page numbers during PDF/PPTX extraction 
   to allow for accurate citations in the RAG chat.
"""

import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def process_single_file(file_obj, kb, filename=None, full_path=None):
    """
    Orchestrates the ingestion of different file types into the KnowledgeBase.
    
    The Lifecycle:
    1. **Format Detection**: Routes by extension.
    2. **Extraction**: Converts binary/unstructured data to raw text.
    3. **Handover**: Sends text to `kb.process_text` for chunking and vectorization.

    Args:
        file_obj: The file-like object (UploadedFile) or a path string.
        kb: The KnowledgeBase instance where the results will be stored.
        filename (str): The display name used for future citations.
        full_path (str): The absolute disk path (crucial for local-file indexing).
        
    Returns:
        str: The name of the processed file for UI reporting.
    """

    fname = filename if filename else getattr(file_obj, 'name', 'unknown_file')
    
    # 1. PDF Handler: Iterates through pages and extracts raw text string.
    if fname.endswith(".pdf"):
        reader = PdfReader(file_obj)
        for i, page in enumerate(reader.pages):
            p_text = page.extract_text() or ""
            # Handover to KB: Text is cleaned and chunked within the KnowledgeBase class.
            kb.process_text(fname, p_text, i + 1, full_path=full_path)
            
    # 2. Tabular Handler (CSV): Converts rows to semantic strings.
    elif fname.endswith(".csv"):
        try:
            df = pd.read_csv(file_obj)
        except UnicodeDecodeError:
            # Fallback for common Windows/European encodings
            if hasattr(file_obj, 'seek'): file_obj.seek(0)
            try:
                df = pd.read_csv(file_obj, encoding='latin1')
            except UnicodeDecodeError:
                if hasattr(file_obj, 'seek'): file_obj.seek(0)
                df = pd.read_csv(file_obj, encoding='cp1252')
        kb.process_dataset(fname, df)
        
    # 3. Tabular Handler (Excel): Supports both .xls and .xlsx formats.
    elif fname.endswith((".xls", ".xlsx")):
        df = pd.read_excel(file_obj)
        kb.process_dataset(fname, df)
        
    # 4. Text/Markdown Handler: Reads full content as a single semantic unit (initially).
    elif fname.endswith((".md", ".txt")):
        encodings = ['utf-8', 'latin1', 'cp1252']
        content = None
        
        if hasattr(file_obj, 'read'):
            raw_data = file_obj.read()
            for enc in encodings:
                try:
                    content = raw_data.decode(enc)
                    break
                except UnicodeDecodeError:
                    continue
            if content is None: raise UnicodeDecodeError("Unable to decode file with utf-8, latin1, or cp1252")
        else:
            for enc in encodings:
                try:
                    with open(file_obj, 'r', encoding=enc) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            if content is None: raise UnicodeDecodeError("Unable to decode file with utf-8, latin1, or cp1252")
            
        kb.process_text(fname, content, 1, full_path=full_path)
        
    # 5. Microsoft Word Handler (.docx)
    elif fname.endswith(".docx"):
        doc = Document(file_obj)
        content = "\n".join([para.text for para in doc.paragraphs])
        kb.process_text(fname, content, 1, full_path=full_path)

    # 6. Microsoft PowerPoint Handler (.pptx)
    elif fname.endswith(".pptx"):
        prs = Presentation(file_obj)
        text_runs = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                full_slide_text = "\n".join(slide_text)
                kb.process_text(fname, full_slide_text, i + 1, full_path=full_path)
        
    return fname
